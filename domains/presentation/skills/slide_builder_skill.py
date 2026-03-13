import os
from pptx.presentation import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from domains.presentation.skills.theme_skill import ThemeSkill


class SlideBuilderSkill:
    """개별 슬라이드 생성 스킬"""

    def __init__(self):
        self._theme = ThemeSkill()

    def build_cover(self, prs: Presentation, product_name: str) -> None:
        """표지 슬라이드 생성"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._theme.apply_background(slide, self._theme.COLOR_TITLE)

        # 상품명
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = product_name
        self._theme.apply_text_style(run, self._theme.FONT_COVER_TITLE, self._theme.COLOR_WHITE, bold=True)

        # 부제목
        sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(0.8))
        tf = sub_box.text_frame
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = "B2B 상품소개서"
        self._theme.apply_text_style(run, self._theme.FONT_COVER_SUBTITLE, self._theme.COLOR_SUBTITLE)

    def build_content(self, prs: Presentation, slide_data: dict) -> None:
        """본문 슬라이드 생성"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 이미지 삽입 (우측)
        if slide_data.get("image_path") and os.path.exists(slide_data["image_path"]):
            slide.shapes.add_picture(
                slide_data["image_path"],
                Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.5),
            )

        # 제목
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.9))
        run = title_box.text_frame.paragraphs[0].add_run()
        run.text = slide_data.get("title", "")
        self._theme.apply_text_style(run, self._theme.FONT_SLIDE_TITLE, self._theme.COLOR_TITLE, bold=True)

        # 액센트 바
        bar = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._theme.COLOR_ACCENT

        # 본문
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.2), Inches(4.5))
        body_box.text_frame.word_wrap = True
        run = body_box.text_frame.paragraphs[0].add_run()
        run.text = slide_data.get("body", "")
        self._theme.apply_text_style(run, self._theme.FONT_SLIDE_BODY, self._theme.COLOR_BODY)

        # 핵심 메시지
        key_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12), Inches(0.8))
        run = key_box.text_frame.paragraphs[0].add_run()
        run.text = f"✦ {slide_data.get('key_message', '')}"
        self._theme.apply_text_style(run, self._theme.FONT_KEY_MESSAGE, self._theme.COLOR_ACCENT, italic=True)

    def build_closing(self, prs: Presentation, product_name: str) -> None:
        """마무리 슬라이드 생성"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._theme.apply_background(slide, self._theme.COLOR_ACCENT)

        text_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10), Inches(1.5))
        tf = text_box.text_frame
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = "감사합니다"
        self._theme.apply_text_style(run, self._theme.FONT_CLOSING, self._theme.COLOR_WHITE, bold=True)
