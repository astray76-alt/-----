import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from shared.config import Config


class PresentationService:
    """
    PPT 생성 서비스 클래스
    - python-pptx를 사용해 슬라이드 파일 생성
    - 콘텐츠 + 디자이너 에이전트 결과물을 PPT로 조립
    """

    # 슬라이드 기본 색상 설정
    _COLOR_TITLE = RGBColor(0x1F, 0x39, 0x64)    # 진한 남색
    _COLOR_BODY = RGBColor(0x33, 0x33, 0x33)      # 진한 회색
    _COLOR_ACCENT = RGBColor(0x2E, 0x75, 0xB6)    # 파란색 강조

    def __init__(self):
        # 출력 경로 생성
        os.makedirs(Config.OUTPUT_DIR, exist_ok=True)

    def create_presentation(self, product_name: str, slides_data: list[dict]) -> str:
        """
        슬라이드 데이터로 PPT 파일 생성

        Args:
            product_name: 상품명 (파일명에 사용)
            slides_data: agent에서 조합한 슬라이드별 데이터 리스트
                         각 항목: {title, body, key_message, image_path}

        Returns:
            생성된 PPT 파일 경로
        """
        prs = Presentation()

        # 슬라이드 크기를 16:9 와이드 설정
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 표지 슬라이드 추가
        self._add_cover_slide(prs, product_name)

        # 본문 슬라이드 추가
        for slide_data in slides_data:
            self._add_content_slide(prs, slide_data)

        # 마무리 슬라이드 추가
        self._add_closing_slide(prs, product_name)

        # 파일 저장
        safe_name = product_name.replace(" ", "_")
        output_path = os.path.join(Config.OUTPUT_DIR, f"{safe_name}_상품소개서.pptx")
        prs.save(output_path)

        return output_path

    def _add_cover_slide(self, prs: Presentation, product_name: str) -> None:
        """표지 슬라이드 생성"""
        layout = prs.slide_layouts[6]  # 빈 레이아웃 사용
        slide = prs.slides.add_slide(layout)

        # 배경색 설정 (진한 남색)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._COLOR_TITLE

        # 상품명 텍스트 박스
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5), Inches(10), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = product_name
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # 부제목 텍스트 박스
        sub_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(4.2), Inches(10), Inches(0.8)
        )
        tf = sub_box.text_frame
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = "B2B 상품소개서"
        run.font.size = Pt(24)
        run.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)

    def _add_content_slide(self, prs: Presentation, slide_data: dict) -> None:
        """본문 슬라이드 생성"""
        layout = prs.slide_layouts[6]  # 빈 레이아웃 사용
        slide = prs.slides.add_slide(layout)

        # 이미지가 있으면 배경에 삽입 (우측 절반)
        if slide_data.get("image_path") and os.path.exists(slide_data["image_path"]):
            slide.shapes.add_picture(
                slide_data["image_path"],
                Inches(7.0), Inches(1.2),
                Inches(5.8), Inches(5.5),
            )

        # 제목 텍스트 박스
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.9)
        )
        tf = title_box.text_frame
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = slide_data.get("title", "")
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = self._COLOR_TITLE

        # 구분선 역할 텍스트 박스 (액센트 컬러 바)
        bar_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.08)
        )
        bar_box.fill.solid()
        bar_box.fill.fore_color.rgb = self._COLOR_ACCENT

        # 본문 텍스트 박스
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(6.2), Inches(4.5)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = slide_data.get("body", "")
        run.font.size = Pt(16)
        run.font.color.rgb = self._COLOR_BODY

        # 핵심 메시지 텍스트 박스 (하단)
        key_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.3), Inches(12), Inches(0.8)
        )
        tf = key_box.text_frame
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = f"✦ {slide_data.get('key_message', '')}"
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.color.rgb = self._COLOR_ACCENT

    def _add_closing_slide(self, prs: Presentation, product_name: str) -> None:
        """마무리 슬라이드 생성"""
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # 배경색 설정
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._COLOR_ACCENT

        # 마무리 문구 텍스트 박스
        text_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.8), Inches(10), Inches(1.5)
        )
        tf = text_box.text_frame
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = f"감사합니다"
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
